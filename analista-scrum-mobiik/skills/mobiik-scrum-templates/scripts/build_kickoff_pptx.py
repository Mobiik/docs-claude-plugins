"""
Kickoff PPT — Migración Preselección TecMilenio
Estándar visual Mobiik: fondo oscuro, acento #AADC1E, Arial Black titulares.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── Brand palette ────────────────────────────────────────────────
BG_BLACK     = RGBColor(0x0A, 0x0A, 0x0A)
BG_PANEL     = RGBColor(0x18, 0x18, 0x18)
BG_PANEL_LT  = RGBColor(0x25, 0x25, 0x25)
LIME         = RGBColor(0xAA, 0xDC, 0x1E)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GREY_TXT     = RGBColor(0xAA, 0xAA, 0xAA)
GREY_DIM     = RGBColor(0x66, 0x66, 0x66)
ORANGE       = RGBColor(0xFF, 0xA0, 0x40)
RED_CRIT     = RGBColor(0xFF, 0x44, 0x44)
BLUE_INFO    = RGBColor(0x1E, 0xA0, 0xDC)

FONT_HEAD = "Arial Black"
FONT_BODY = "Calibri"

# ─── Slide setup (widescreen 16:9) ────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]


# ─── Helpers ──────────────────────────────────────────────────────
def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_rect(slide, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def add_text(slide, x, y, w, h, text, *, font=FONT_BODY, size=14,
             color=WHITE, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(line, tuple):
            for chunk_text, chunk_opts in line:
                r = p.add_run()
                r.text = chunk_text
                r.font.name = chunk_opts.get("font", font)
                r.font.size = Pt(chunk_opts.get("size", size))
                r.font.bold = chunk_opts.get("bold", bold)
                r.font.italic = chunk_opts.get("italic", italic)
                r.font.color.rgb = chunk_opts.get("color", color)
        else:
            r = p.add_run()
            r.text = line
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
    return tb

def add_bg(slide, color=BG_BLACK):
    add_rect(slide, 0, 0, SW, SH, color)

def add_footer(slide):
    add_text(slide, Inches(0.4), Inches(7.05), Inches(7), Inches(0.3),
             "© 2026 Mobiik  |  coding for the future",
             font=FONT_BODY, size=9, color=GREY_DIM)

def add_section_num(slide, num, x=Inches(11.0), y=Inches(0.35)):
    if not num:
        return
    tb = slide.shapes.add_textbox(x, y, Inches(2.2), Inches(1.4))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = num
    r.font.name = FONT_HEAD
    r.font.size = Pt(64)
    r.font.bold = True
    r.font.color.rgb = LIME
    return tb

def add_title(slide, title, *, y=Inches(0.45), size=32):
    add_text(slide, Inches(0.5), y, Inches(11), Inches(0.7), title,
             font=FONT_HEAD, size=size, color=WHITE, bold=True)
    # Lime underbar (NOT a line — a small block accent)
    bar = add_rect(slide, Inches(0.5), y + Inches(0.65), Inches(0.6), Emu(45000), LIME)

def add_accent_card(slide, x, y, w, h, *, fill_color=BG_PANEL):
    """Card with thin left lime accent bar."""
    add_rect(slide, x, y, w, h, fill_color)
    add_rect(slide, x, y, Emu(60000), h, LIME)

def add_image(slide, path, x, y, w, h):
    return slide.shapes.add_picture(path, x, y, width=w, height=h)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s, BG_BLACK)

# Cosmic image as right-side full bleed
add_image(s, "./assets/cosmic_dim.jpeg",
          Inches(7.3), Inches(0), Inches(6.03), SH)
# Dark gradient overlay on left half for text legibility
add_rect(s, 0, 0, Inches(8.5), SH, BG_BLACK)

# KICK-OFF eyebrow
add_text(s, Inches(0.7), Inches(0.9), Inches(6), Inches(0.4),
         "KICK-OFF", font=FONT_HEAD, size=14, color=LIME, bold=True)

# Title
add_text(s, Inches(0.7), Inches(1.35), Inches(8), Inches(1.2),
         "MIGRACIÓN", font=FONT_HEAD, size=52, color=WHITE, bold=True)
add_text(s, Inches(0.7), Inches(2.20), Inches(8), Inches(1.2),
         "PRESELECCIÓN", font=FONT_HEAD, size=52, color=LIME, bold=True)

# Subtitle
add_text(s, Inches(0.7), Inches(3.2), Inches(8), Inches(0.5),
         "Tecmilenio  ×  Mobiik   ·   Lift-and-Shift a Azure PaaS",
         font=FONT_BODY, size=15, color=GREY_TXT, italic=True)

# Stat row — 3 cards
card_y = Inches(4.1)
card_h = Inches(1.7)
card_w = Inches(2.5)
gap    = Inches(0.15)
for i, (label, value) in enumerate([
    ("FECHA",     "Junio 2026"),
    ("INVERSIÓN", "$218,790 MXN"),
    ("DURACIÓN",  "3 sem · 15 días"),
]):
    x = Inches(0.7) + (card_w + gap) * i
    add_accent_card(s, x, card_y, card_w, card_h, fill_color=BG_PANEL)
    add_text(s, x + Inches(0.25), card_y + Inches(0.25),
             card_w - Inches(0.4), Inches(0.3),
             label, font=FONT_HEAD, size=10, color=LIME, bold=True)
    add_text(s, x + Inches(0.25), card_y + Inches(0.65),
             card_w - Inches(0.4), Inches(0.9),
             value, font=FONT_HEAD, size=18, color=WHITE, bold=True)

add_footer(s)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_title(s, "AGENDA")
add_section_num(s, "")  # no number on agenda

agenda = [
    ("01", "Contexto y Objetivo",     "Migración técnica de Preselección a Azure PaaS"),
    ("02", "Alcance del Proyecto",    "Lift-and-Shift, hallazgos de seguridad y BD"),
    ("03", "Arquitectura Azure",      "Componentes PaaS objetivo de la migración"),
    ("04", "Plan de Trabajo",         "Cronograma 3 semanas / 15 días útiles"),
    ("05", "Equipo y Dependencias",   "Roles Mobiik y dependencias con Tecmilenio"),
    ("06", "Inversión y Próximos Pasos", "Esquema de pagos, entregables y siguientes acciones"),
]

col_w = Inches(6.0)
row_h = Inches(1.45)
for idx, (num, title, desc) in enumerate(agenda):
    col = idx % 2
    row = idx // 2
    x = Inches(0.5) + col * (col_w + Inches(0.3))
    y = Inches(1.75) + row * (row_h + Inches(0.2))
    add_accent_card(s, x, y, col_w, row_h)
    add_text(s, x + Inches(0.3), y + Inches(0.2), Inches(0.9), Inches(1),
             num, font=FONT_HEAD, size=34, color=LIME, bold=True)
    add_text(s, x + Inches(1.4), y + Inches(0.25), col_w - Inches(1.6), Inches(0.5),
             title, font=FONT_HEAD, size=18, color=WHITE, bold=True)
    add_text(s, x + Inches(1.4), y + Inches(0.75), col_w - Inches(1.6), Inches(0.6),
             desc, font=FONT_BODY, size=12, color=GREY_TXT)

add_footer(s)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — CONTEXTO Y OBJETIVO
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_title(s, "CONTEXTO Y OBJETIVO DEL PROYECTO")
add_section_num(s, "01")

# Description block
add_text(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(1.3),
         "Tecmilenio requiere migrar técnicamente el aplicativo institucional "
         "Preselección hacia una infraestructura PaaS en Microsoft Azure. El proyecto "
         "se ejecuta bajo el modelo Lift-and-Shift: se traslada la solución actual "
         "IaaS a un entorno PaaS administrado sin intervenir la lógica de negocio "
         "ni realizar reingeniería del aplicativo.",
         font=FONT_BODY, size=13, color=GREY_TXT)

# 4 pillar cards
pillars = [
    ("📁", "Continuidad Operativa",
     "Garantizar la operación de Preselección durante y después de la migración a Azure"),
    ("🔍", "Lift-and-Shift",
     "Migración sin modificar lógica de negocio ni reescritura del aplicativo"),
    ("☁️", "Azure PaaS",
     "Trasladar la solución actual IaaS a un entorno PaaS administrado de Microsoft"),
    ("🔐", "Seguridad y Calidad",
     "Atención de los 13 hallazgos de seguridad identificados por Tecmilenio"),
]
card_w = Inches(2.95); card_h = Inches(2.8); gap = Inches(0.15)
for i, (icon, title, desc) in enumerate(pillars):
    x = Inches(0.5) + i * (card_w + gap)
    y = Inches(3.6)
    add_accent_card(s, x, y, card_w, card_h)
    add_text(s, x + Inches(0.25), y + Inches(0.25), card_w - Inches(0.4), Inches(0.7),
             icon, font=FONT_BODY, size=32, color=LIME)
    add_text(s, x + Inches(0.25), y + Inches(1.05), card_w - Inches(0.4), Inches(0.6),
             title, font=FONT_HEAD, size=15, color=WHITE, bold=True)
    add_text(s, x + Inches(0.25), y + Inches(1.6), card_w - Inches(0.4), Inches(1.1),
             desc, font=FONT_BODY, size=11, color=GREY_TXT)

add_footer(s)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — ALCANCE DEL PROYECTO
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_title(s, "ALCANCE DEL PROYECTO")
add_section_num(s, "02")

scope_blocks = [
    ("Análisis Técnico", [
        "Revisión arquitectura actual",
        "Identificación de dependencias",
        "Validación requerimientos Azure",
    ]),
    ("Preparación de Entornos", [
        "Resource Groups y App Services",
        "Key Vault, Storage, Redis",
        "Ambientes No-Prod y Prod",
    ]),
    ("Migración Lift-and-Shift", [
        "Despliegue en App Service",
        "Variables, secretos, conexiones",
        "Sin cambios a lógica de negocio",
    ]),
    ("Atención de Hallazgos", [
        "13 hallazgos de seguridad",
        "Actualización de librerías",
        "Validación y reescaneo",
    ]),
    ("Bases de Datos y DevOps", [
        "Migración esquemas MS SQL",
        "Pipelines CI/CD",
        "Pruebas smoke con datos",
    ]),
    ("Pruebas y Go-Live", [
        "Smoke tests e integraciones",
        "Acompañamiento a UAT",
        "Despliegue productivo",
    ]),
]

cw = Inches(4.05); ch = Inches(2.45); gx = Inches(0.15); gy = Inches(0.2)
for i, (title, items) in enumerate(scope_blocks):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + col * (cw + gx)
    y = Inches(1.85) + row * (ch + gy)
    add_accent_card(s, x, y, cw, ch)
    add_text(s, x + Inches(0.25), y + Inches(0.2), cw - Inches(0.4), Inches(0.45),
             title, font=FONT_HEAD, size=14, color=LIME, bold=True)
    bullet_lines = [(("●  ", {"color": LIME, "size": 11, "bold": True}),
                     (it, {"color": WHITE, "size": 11})) for it in items]
    add_text(s, x + Inches(0.25), y + Inches(0.75), cw - Inches(0.4), Inches(1.6),
             bullet_lines, font=FONT_BODY, size=11, color=WHITE)

add_footer(s)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — ARQUITECTURA AZURE PAAS
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_title(s, "ARQUITECTURA AZURE PAAS")
add_section_num(s, "03")

arch = [
    ("App Service",      "Backend .NET Framework + MVC migrado a Azure App Service administrado"),
    ("Static Web Apps",  "Frontend (AngularJS) desplegado en Static Web Apps cuando aplique"),
    ("Key Vault",        "Gestión centralizada de secretos, cadenas de conexión y credenciales"),
    ("Storage / Redis",  "Storage y Azure Managed Redis en caso de ser requeridos por el aplicativo"),
    ("Azure SQL",        "Migración de esquemas y datos MS SQL Server desde IaaS hacia Azure PaaS"),
    ("Azure DevOps",     "Pipelines de CI/CD integrados con repositorios existentes hacia los ambientes"),
]

cw = Inches(4.05); ch = Inches(1.7); gx = Inches(0.15); gy = Inches(0.2)
for i, (title, desc) in enumerate(arch):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + col * (cw + gx)
    y = Inches(1.85) + row * (ch + gy)
    add_accent_card(s, x, y, cw, ch)
    add_text(s, x + Inches(0.25), y + Inches(0.2), cw - Inches(0.4), Inches(0.45),
             title, font=FONT_HEAD, size=15, color=WHITE, bold=True)
    add_text(s, x + Inches(0.25), y + Inches(0.7), cw - Inches(0.4), Inches(0.9),
             desc, font=FONT_BODY, size=11, color=GREY_TXT)

# Bottom callout
callout_y = Inches(5.55)
add_rect(s, Inches(0.5), callout_y, Inches(12.3), Inches(1.1), BG_PANEL_LT)
add_rect(s, Inches(0.5), callout_y, Emu(60000), Inches(1.1), LIME)
add_text(s, Inches(0.85), callout_y + Inches(0.18), Inches(11.8), Inches(0.4),
         "💡  Lift-and-Shift",
         font=FONT_HEAD, size=13, color=LIME, bold=True)
add_text(s, Inches(0.85), callout_y + Inches(0.55), Inches(11.8), Inches(0.5),
         "Se conserva la familia tecnológica actual (.NET Framework 4.7.x → 4.8.1, "
         "AngularJS 1.5 → 1.8.x). Sin reescritura ni migración a nuevos frameworks.",
         font=FONT_BODY, size=11, color=WHITE)

add_footer(s)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — PLAN DE TRABAJO
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_title(s, "PLAN DE TRABAJO")
add_section_num(s, "04")

sprints = [
    ("Sprint 1", "Semana 1 · L–M (2 días)", [
        "Análisis técnico",
        "Arquitectura objetivo",
        "Dependencias",
        "Plan detallado",
    ]),
    ("Sprint 2", "Semana 1–2 · X–L (4 días)", [
        "Lift-and-Shift",
        "App Service / WebApp",
        "Configuración Azure",
        "DB MS SQL Server",
    ]),
    ("Sprint 3", "Semana 2 · M–V (4 días)", [
        "13 Hallazgos seg.",
        "QA interno",
        "Smoke tests",
        "Pipelines CI/CD",
    ]),
    ("UAT", "Semana 3 · L–X (3 días)", [
        "Pruebas cliente",
        "Corrección incidentes",
        "Validación final",
        "Aceptación",
    ]),
    ("Go-Live", "Semana 3 · Jueves (1 día)", [
        "Despliegue PROD",
        "Acompañamiento",
        "Carta de cierre",
        "Hand-off",
    ]),
]

cw = Inches(2.45); ch = Inches(3.5); gx = Inches(0.1)
total_w = cw * 5 + gx * 4
start_x = (SW - total_w) / 2
for i, (name, when, items) in enumerate(sprints):
    x = start_x + i * (cw + gx)
    y = Inches(1.9)
    add_accent_card(s, x, y, cw, ch)
    # Sprint name big
    add_text(s, x + Inches(0.2), y + Inches(0.2), cw - Inches(0.3), Inches(0.5),
             name, font=FONT_HEAD, size=18, color=LIME, bold=True)
    # When
    add_text(s, x + Inches(0.2), y + Inches(0.7), cw - Inches(0.3), Inches(0.6),
             when, font=FONT_BODY, size=10, color=GREY_TXT, italic=True)
    # Separator
    add_rect(s, x + Inches(0.2), y + Inches(1.25), cw - Inches(0.4), Emu(15000), LIME)
    # Items
    lines = [(("● ", {"color": LIME, "size": 10, "bold": True}),
              (it, {"color": WHITE, "size": 11})) for it in items]
    add_text(s, x + Inches(0.2), y + Inches(1.4), cw - Inches(0.3), Inches(2),
             lines, font=FONT_BODY, size=11, color=WHITE)

# Warning callout
warn_y = Inches(5.65)
add_rect(s, Inches(0.5), warn_y, Inches(12.3), Inches(0.95), BG_PANEL_LT)
add_rect(s, Inches(0.5), warn_y, Emu(60000), Inches(0.95), ORANGE)
add_text(s, Inches(0.85), warn_y + Inches(0.2), Inches(11.8), Inches(0.55),
         [(("⚠️  ", {"size": 14, "color": ORANGE}),
           ("El cliente debe entregar accesos a Azure, repositorios GitHub y documentación técnica en los primeros 2 días del proyecto.",
            {"size": 12, "color": WHITE, "bold": True}))])

add_footer(s)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — DEPENDENCIAS TECMILENIO
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_title(s, "DEPENDENCIAS TECMILENIO")
add_section_num(s, "05")

deps = [
    ("CRÍTICO",      RED_CRIT,   "Sprint 1",    "Accesos a repositorios GitHub, suscripción Azure y permisos para crear recursos (App Service, Key Vault, SQL)"),
    ("CRÍTICO",      RED_CRIT,   "Sprint 1",    "Documentación técnica actual del aplicativo Preselección: arquitectura, integraciones, pipelines y dependencias"),
    ("ALTA",         ORANGE,     "Sprint 1-2",  "Encargado de Infraestructura disponible para configuración de Azure, repositorios GitHub y permisos"),
    ("ALTA",         ORANGE,     "Sprint 2",    "Confirmación formal de los 13 hallazgos de seguridad incluidos en el alcance (sin adicionales durante ejecución)"),
    ("MEDIA",        LIME,       "Sprint 3",    "Equipo de Pruebas del cliente disponible para UAT durante los 3 días asignados (Semana 3, L-X)"),
    ("MEDIA",        LIME,       "Sprint 3",    "Provisión y configuración de la infraestructura productiva en Azure por parte del equipo Tecmilenio"),
    ("INFORMATIVA",  BLUE_INFO,  "Continuo",    "Product Owner disponible para validación funcional y toma de decisiones durante todo el proyecto"),
    ("INFORMATIVA",  BLUE_INFO,  "Cierre",      "Validación funcional completa del aplicativo en producción posterior a la liberación es responsabilidad del cliente"),
]

row_h = Inches(0.6); gy = Inches(0.05)
for i, (severity, color, when, text) in enumerate(deps):
    y = Inches(1.85) + i * (row_h + gy)
    add_rect(s, Inches(0.5), y, Inches(12.3), row_h, BG_PANEL)
    # Severity pill
    add_rect(s, Inches(0.65), y + Inches(0.13), Inches(1.5), Inches(0.35), color)
    add_text(s, Inches(0.65), y + Inches(0.16), Inches(1.5), Inches(0.3),
             severity, font=FONT_HEAD, size=9, color=BG_BLACK, bold=True,
             align=PP_ALIGN.CENTER)
    # When
    add_text(s, Inches(2.35), y + Inches(0.16), Inches(1.3), Inches(0.3),
             when, font=FONT_HEAD, size=10, color=GREY_TXT, bold=True)
    # Description
    add_text(s, Inches(3.75), y + Inches(0.14), Inches(8.9), Inches(0.4),
             text, font=FONT_BODY, size=11, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

add_footer(s)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — EQUIPO DEL PROYECTO
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_title(s, "EQUIPO DEL PROYECTO")
add_section_num(s, "06")

# MOBIIK header
add_text(s, Inches(0.5), Inches(1.85), Inches(6), Inches(0.4),
         "MOBIIK", font=FONT_HEAD, size=14, color=LIME, bold=True)

team = [
    ("Miroslava Jiménez", "Scrum Master"),
    ("Andrés Hernández",  "Arquitecto"),
    ("Iran Arellanes",    "Dev Fullstack"),
    ("Norberto Galicia",  "DevOps"),
    ("César Zúñiga",      "QA Tester"),
]
cw = Inches(2.3); ch = Inches(1.45); gx = Inches(0.1)
for i, (name, role) in enumerate(team):
    x = Inches(0.5) + i * (cw + gx)
    y = Inches(2.35)
    add_accent_card(s, x, y, cw, ch)
    add_text(s, x + Inches(0.25), y + Inches(0.2), cw - Inches(0.4), Inches(0.5),
             role, font=FONT_HEAD, size=10, color=LIME, bold=True)
    add_text(s, x + Inches(0.25), y + Inches(0.6), cw - Inches(0.4), Inches(0.8),
             name, font=FONT_HEAD, size=13, color=WHITE, bold=True)

# Modelo de Gobernanza panel
gov_y = Inches(4.15)
add_accent_card(s, Inches(0.5), gov_y, Inches(12.3), Inches(2.7))
add_text(s, Inches(0.85), gov_y + Inches(0.25), Inches(11.8), Inches(0.5),
         "Modelo de Gobernanza",
         font=FONT_HEAD, size=16, color=LIME, bold=True)
gov_items = [
    "Metodología ágil (Scrum) — Sprints semanales con visibilidad continua",
    "Reuniones diarias (Daily Scrum) y revisiones al cierre de cada Sprint",
    "Canal dedicado en Teams para comunicación operativa del proyecto",
    "Escalación:  Scrum Master Mobiik  ↔  Product Owner Tecmilenio",
    "Change orders documentados con impacto en cronograma y costo",
]
lines = [(("●  ", {"color": LIME, "size": 12, "bold": True}),
          (it, {"color": WHITE, "size": 12})) for it in gov_items]
add_text(s, Inches(0.85), gov_y + Inches(0.85), Inches(11.5), Inches(1.7),
         lines, font=FONT_BODY, size=12)

add_footer(s)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — INVERSIÓN Y ENTREGABLES
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_title(s, "INVERSIÓN Y ENTREGABLES")
add_section_num(s, "")

# Left side: investment stats
stats = [
    ("Inversión Total",    "$218,790",   "MXN — Antes de IVA"),
    ("Pago 1 · Kickoff",   "30%",        "$65,637 MXN — Inicio Sprint 1"),
    ("Pago 2 · Cierre",    "70%",        "$153,153 MXN — Cierre Sprint 3"),
    ("Duración",           "3 sem",      "15 días útiles — 3 Sprints"),
]
sx = Inches(0.5); sy_start = Inches(1.85)
sh_card = Inches(1.18); sgap = Inches(0.1)
for i, (label, big, small) in enumerate(stats):
    y = sy_start + i * (sh_card + sgap)
    add_accent_card(s, sx, y, Inches(5.8), sh_card)
    add_text(s, sx + Inches(0.25), y + Inches(0.15), Inches(3.5), Inches(0.4),
             label, font=FONT_HEAD, size=11, color=GREY_TXT, bold=True)
    add_text(s, sx + Inches(0.25), y + Inches(0.45), Inches(2.5), Inches(0.7),
             big, font=FONT_HEAD, size=24, color=LIME, bold=True)
    add_text(s, sx + Inches(2.9), y + Inches(0.55), Inches(2.8), Inches(0.5),
             small, font=FONT_BODY, size=10, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

# Right side: deliverables
dx = Inches(6.6); dy = Inches(1.85)
add_accent_card(s, dx, dy, Inches(6.2), Inches(5.05))
add_text(s, dx + Inches(0.3), dy + Inches(0.2), Inches(5.8), Inches(0.5),
         "ENTREGABLES DEL PROYECTO",
         font=FONT_HEAD, size=14, color=LIME, bold=True)
delivers = [
    "Código fuente de la aplicación migrada",
    "Infraestructura Azure (ambiente no productivo)",
    "Análisis técnico inicial de la aplicación",
    "Memoria técnica de cambios aplicados",
    "Pipelines de CI/CD configurados",
    "Documentación de CI/CD e infraestructura",
    "Corrección de 13 hallazgos de seguridad",
    "Resultados de smoke test",
    "Carta de cierre del proyecto",
]
del_lines = [(("✓  ", {"color": LIME, "size": 13, "bold": True}),
              (d, {"color": WHITE, "size": 12})) for d in delivers]
add_text(s, dx + Inches(0.3), dy + Inches(0.8), Inches(5.8), Inches(4),
         del_lines, font=FONT_BODY, size=12)

add_footer(s)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — PRÓXIMOS PASOS
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_title(s, "PRÓXIMOS PASOS")
add_section_num(s, "")

steps = [
    ("1", "HOY",        "Firma de propuesta",            "Firma de propuesta comercial y orden de inicio del proyecto"),
    ("2", "KICKOFF",    "MAY 22",                        "Kickoff oficial — alineación de equipos Mobiik × Tecmilenio y arranque Sprint 1"),
    ("3", "ACCESOS",    "MAY 22-26",                     "Acceso a repositorios GitHub, infraestructura Azure actual y documentación técnica"),
    ("4", "SPRINT 1",   "MAY 25-29",                     "Sprint 1: Análisis técnico y validación de arquitectura objetivo en Azure PaaS"),
    ("5", "SPRINT 2-3", "JUN 1-12",                      "Sprint 2-3: Migración Lift-and-Shift, atención de 13 hallazgos, smoke tests y UAT"),
    ("6", "GO-LIVE",    "JUN 11",                        "Go-Live productivo y acompañamiento post-liberación inicial"),
]

row_h = Inches(0.78); gy = Inches(0.08)
for i, (num, milestone, when, desc) in enumerate(steps):
    y = Inches(1.85) + i * (row_h + gy)
    add_rect(s, Inches(0.5), y, Inches(12.3), row_h, BG_PANEL)
    # Big number on left
    add_rect(s, Inches(0.5), y, Inches(0.85), row_h, BG_PANEL_LT)
    add_text(s, Inches(0.5), y, Inches(0.85), row_h,
             num, font=FONT_HEAD, size=26, color=LIME, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Milestone tag
    add_text(s, Inches(1.55), y + Inches(0.18), Inches(2.0), Inches(0.42),
             milestone, font=FONT_HEAD, size=13, color=LIME, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    # When
    add_text(s, Inches(3.65), y + Inches(0.18), Inches(1.7), Inches(0.42),
             when, font=FONT_HEAD, size=10, color=GREY_TXT, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    # Description
    add_text(s, Inches(5.45), y + Inches(0.2), Inches(7.2), Inches(0.5),
             desc, font=FONT_BODY, size=11, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

add_footer(s)


# ═══════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════
out = "./entregables/KO TECMILENIO 2.6 - Mobiik Standard.pptx"
prs.save(out)
print(f"✓ Saved: {out}")
print(f"  Slides: {len(prs.slides)}")
