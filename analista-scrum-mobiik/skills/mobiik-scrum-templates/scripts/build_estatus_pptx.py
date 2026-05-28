"""
Presentación de Estatus de Proyecto — Mobiik
Reporte ejecutivo de avance (datos provenientes de Jira), 4 slides.
Estándar visual Mobiik: fondo oscuro #0A0A0A, acento verde lima #AADC1E,
logo MOB·IIK, footer "Mobiik — AI, Cloud & Software Development".

Basado en la plantilla corporativa Mobiik_Estatus_Proyecto_PLANTILLA.pptx
Dependencia:  pip3 install --user python-pptx
Uso:          python3 build_estatus_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── DATOS DEL REPORTE (editar — alimentar con datos de Jira) ────────
PROYECTO   = "[Proyecto]"
PERIODO    = "[Sprint / semana]"
PM         = "[Nombre]"
# Salud general: "VERDE" (en tiempo) | "AMBAR" (en riesgo) | "ROJO" (desviado)
SALUD      = "VERDE"

RESUMEN    = "[2–3 líneas sobre el avance global, lo más relevante del periodo y la perspectiva.]"

# Métricas Jira
M_TOTAL       = "[#]"
M_COMPLETADOS = "[#]"
M_PROGRESO    = "[#]"
M_PORHACER    = "[#]"
M_BLOQUEADOS  = "[#]"
AVANCE_GLOBAL = "[__%]"
VELOCIDAD     = "[__ pts]"
SPRINT_ACTUAL = "[Sprint X]"

LOGROS = [
    "[Issue completado / hito alcanzado — clave Jira]",
    "[...]",
    "[...]",
]

# (Clave, Resumen, Responsable, Estado)
WIP = [
    ("[Clave]", "[Resumen del issue]", "[Responsable]", "[Estado]"),
]

# (Descripción, Mitigación / acción, Responsable, Severidad)
RIESGOS = [
    ("[Descripción del riesgo]", "[Acción de mitigación]", "[Resp.]", "Media"),
]

PROXIMOS = [
    "[Prioridad para el siguiente periodo]",
    "[...]",
]

DECISIONES = "[Indicar qué se necesita del cliente para no frenar el avance, o “Ninguno”.]"

OUT = "./entregables/Estatus - <proyecto> - <periodo>.pptx"
# ─────────────────────────────────────────────────────────────────────

# ─── Brand palette ───────────────────────────────────────────────────
BG_BLACK     = RGBColor(0x0A, 0x0A, 0x0A)
BG_PANEL     = RGBColor(0x18, 0x18, 0x18)
BG_PANEL_LT  = RGBColor(0x25, 0x25, 0x25)
LIME         = RGBColor(0xAA, 0xDC, 0x1E)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GREY_TXT     = RGBColor(0xAA, 0xAA, 0xAA)
GREY_DIM     = RGBColor(0x66, 0x66, 0x66)
ORANGE       = RGBColor(0xFF, 0xA0, 0x40)
RED_CRIT     = RGBColor(0xFF, 0x44, 0x44)
BLUE_INFO    = RGBColor(0x1E, 0xA0, 0xDC)

FONT_HEAD = "Arial Black"
FONT_BODY = "Calibri"

HEALTH = {
    "VERDE": (LIME,     "● VERDE"),
    "AMBAR": (ORANGE,   "● ÁMBAR"),
    "ROJO":  (RED_CRIT, "● ROJO"),
}
SEV_COLOR = {"alta": RED_CRIT, "media": ORANGE, "baja": LIME}

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


# ─── Helpers ──────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def add_text(slide, x, y, w, h, text, *, font=FONT_BODY, size=14,
             color=WHITE, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def add_bg(slide, color=BG_BLACK):
    add_rect(slide, 0, 0, SW, SH, color)

def add_logo(slide, x=Inches(0.5), y=Inches(0.35)):
    """Logo MOB·IIK — 'MOB' blanco + 'IIK' verde lima."""
    tb = slide.shapes.add_textbox(x, y, Inches(2.0), Inches(0.5))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "MOB"
    r1.font.name = FONT_HEAD; r1.font.size = Pt(20); r1.font.bold = True; r1.font.color.rgb = WHITE
    r2 = p.add_run(); r2.text = "IIK"
    r2.font.name = FONT_HEAD; r2.font.size = Pt(20); r2.font.bold = True; r2.font.color.rgb = LIME
    return tb

def add_footer(slide, page):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(9), Inches(0.3),
             "Mobiik — AI, Cloud & Software Development",
             font=FONT_BODY, size=9, color=GREY_DIM)
    add_text(slide, Inches(12.4), Inches(7.05), Inches(0.6), Inches(0.3),
             str(page), font=FONT_BODY, size=9, color=GREY_DIM, align=PP_ALIGN.RIGHT)

def add_title(slide, title, eyebrow=None):
    if eyebrow:
        add_text(slide, Inches(0.5), Inches(1.05), Inches(11), Inches(0.35),
                 eyebrow, font=FONT_BODY, size=11, color=GREY_TXT, italic=True)
    add_text(slide, Inches(0.5), Inches(1.4), Inches(11.5), Inches(0.7), title,
             font=FONT_HEAD, size=26, color=WHITE, bold=True)
    add_rect(slide, Inches(0.5), Inches(2.0), Inches(0.6), Emu(45000), LIME)

def add_accent_card(slide, x, y, w, h, *, fill_color=BG_PANEL):
    add_rect(slide, x, y, w, h, fill_color)
    add_rect(slide, x, y, Emu(60000), h, LIME)

def add_branded_table(slide, x, y, w, headers, rows, col_fr, row_h=Inches(0.42),
                      hdr_h=Inches(0.42), font_size=11):
    """Tabla con header oscuro y filas alternas, ancho proporcional col_fr."""
    total = sum(col_fr)
    widths = [int(w * fr / total) for fr in col_fr]
    # Header
    cx = x
    for i, h in enumerate(headers):
        add_rect(slide, cx, y, widths[i], hdr_h, BG_PANEL_LT)
        add_text(slide, cx + Inches(0.1), y, widths[i] - Inches(0.2), hdr_h, h,
                 font=FONT_BODY, size=font_size, color=LIME, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        cx += widths[i]
    # Rows
    cy = y + hdr_h
    for ri, row in enumerate(rows):
        bg = BG_PANEL if ri % 2 == 0 else BG_BLACK
        cx = x
        for ci, val in enumerate(row):
            add_rect(slide, cx, cy, widths[ci], row_h, bg)
            col = WHITE
            if headers[ci].lower().startswith("sever") and str(val).lower() in SEV_COLOR:
                col = SEV_COLOR[str(val).lower()]
            add_text(slide, cx + Inches(0.1), cy, widths[ci] - Inches(0.2), row_h,
                     str(val), font=FONT_BODY, size=font_size, color=col,
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += widths[ci]
        cy += row_h


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_logo(s)

add_text(s, Inches(0.8), Inches(2.3), Inches(11), Inches(0.9),
         "Estatus de Proyecto", font=FONT_HEAD, size=44, color=WHITE, bold=True)
add_rect(s, Inches(0.85), Inches(3.25), Inches(0.8), Emu(50000), LIME)
add_text(s, Inches(0.85), Inches(3.45), Inches(11), Inches(0.4),
         "Reporte de avance — datos provenientes de Jira",
         font=FONT_BODY, size=14, color=GREY_TXT, italic=True)

# Cards de datos clave
cards = [("PROYECTO", PROYECTO), ("PERIODO", PERIODO), ("PM", PM)]
cw, ch, gap = Inches(3.0), Inches(1.3), Inches(0.25)
cx, cy = Inches(0.85), Inches(4.4)
for label, value in cards:
    add_accent_card(s, cx, cy, cw, ch)
    add_text(s, cx + Inches(0.25), cy + Inches(0.2), cw - Inches(0.4), Inches(0.3),
             label, font=FONT_BODY, size=11, color=LIME, bold=True)
    add_text(s, cx + Inches(0.25), cy + Inches(0.55), cw - Inches(0.4), Inches(0.6),
             value, font=FONT_HEAD, size=16, color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    cx += cw + gap

# Card de salud general (semáforo)
health_color, health_label = HEALTH.get(SALUD.upper(), HEALTH["VERDE"])
add_rect(s, cx, cy, cw, ch, BG_PANEL)
add_rect(s, cx, cy, Emu(60000), ch, health_color)
add_text(s, cx + Inches(0.25), cy + Inches(0.2), cw - Inches(0.4), Inches(0.3),
         "SALUD GENERAL", font=FONT_BODY, size=11, color=GREY_TXT, bold=True)
add_text(s, cx + Inches(0.25), cy + Inches(0.55), cw - Inches(0.4), Inches(0.6),
         health_label, font=FONT_HEAD, size=18, color=health_color, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.85), Inches(6.1), Inches(11), Inches(0.4),
         "Convención de salud:  ● Verde en tiempo   ● Ámbar en riesgo   ● Rojo desviado",
         font=FONT_BODY, size=10, color=GREY_DIM, italic=True)
add_footer(s, 1)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — RESUMEN EJECUTIVO Y MÉTRICAS
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_logo(s)
add_title(s, "Resumen ejecutivo y métricas", eyebrow="Datos provenientes de Jira")

# Resumen
add_text(s, Inches(0.5), Inches(2.25), Inches(2.0), Inches(0.3),
         "RESUMEN", font=FONT_BODY, size=12, color=LIME, bold=True)
add_accent_card(s, Inches(0.5), Inches(2.6), Inches(12.3), Inches(0.95))
add_text(s, Inches(0.75), Inches(2.7), Inches(11.8), Inches(0.8), RESUMEN,
         font=FONT_BODY, size=13, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# 5 métricas
metrics = [("Total", M_TOTAL, WHITE), ("Completados", M_COMPLETADOS, LIME),
           ("En progreso", M_PROGRESO, BLUE_INFO), ("Por hacer", M_PORHACER, GREY_TXT),
           ("Bloqueados", M_BLOQUEADOS, RED_CRIT)]
mw, mh, mgap = Inches(2.3), Inches(1.25), Inches(0.18)
mx, my = Inches(0.5), Inches(3.85)
for label, value, col in metrics:
    add_rect(s, mx, my, mw, mh, BG_PANEL)
    add_text(s, mx, my + Inches(0.18), mw, Inches(0.6), value,
             font=FONT_HEAD, size=34, color=col, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, mx, my + Inches(0.85), mw, Inches(0.3), label,
             font=FONT_BODY, size=11, color=GREY_TXT, align=PP_ALIGN.CENTER)
    mx += mw + mgap

# Barra de KPIs
kpi_y = Inches(5.35)
add_accent_card(s, Inches(0.5), kpi_y, Inches(12.3), Inches(0.6))
kpis = [("AVANCE GLOBAL:", AVANCE_GLOBAL), ("VELOCIDAD SPRINT:", VELOCIDAD),
        ("SPRINT ACTUAL:", SPRINT_ACTUAL)]
kx = Inches(0.85)
for label, value in kpis:
    add_text(s, kx, kpi_y, Inches(2.0), Inches(0.6), label,
             font=FONT_BODY, size=11, color=GREY_TXT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, kx + Inches(1.85), kpi_y, Inches(1.8), Inches(0.6), value,
             font=FONT_HEAD, size=14, color=LIME, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    kx += Inches(4.0)

# Logros del periodo
add_text(s, Inches(0.5), Inches(6.15), Inches(6), Inches(0.3),
         "Logros del periodo", font=FONT_HEAD, size=14, color=WHITE, bold=True)
ly = Inches(6.5)
for logro in LOGROS:
    add_text(s, Inches(0.6), ly, Inches(0.3), Inches(0.3), "▸",
             font=FONT_BODY, size=11, color=LIME, bold=True)
    add_text(s, Inches(0.95), ly, Inches(11.5), Inches(0.3), logro,
             font=FONT_BODY, size=11, color=WHITE)
    ly += Inches(0.32)
add_footer(s, 2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — DISTRIBUCIÓN Y TRABAJO EN PROGRESO
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_logo(s)
add_title(s, "Distribución de issues y trabajo en progreso",
          eyebrow="Datos provenientes de Jira")

# Placeholder de gráfica de distribución (izquierda)
add_accent_card(s, Inches(0.5), Inches(2.35), Inches(3.6), Inches(4.0))
add_text(s, Inches(0.7), Inches(3.9), Inches(3.2), Inches(1.0),
         ["Gráfica de", "distribución", "(se llena con", "datos de Jira)"],
         font=FONT_BODY, size=12, color=GREY_TXT, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

# Tabla Trabajo en progreso (derecha)
add_text(s, Inches(4.4), Inches(2.35), Inches(8), Inches(0.35),
         "Trabajo en progreso", font=FONT_HEAD, size=14, color=WHITE, bold=True)
add_branded_table(
    s, Inches(4.4), Inches(2.8), Inches(8.4),
    ["Clave", "Resumen", "Responsable", "Estado"], WIP,
    col_fr=[1.2, 3.5, 1.6, 1.4],
)
add_footer(s, 3)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — RIESGOS, PRÓXIMOS PASOS Y DECISIONES
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
add_logo(s)
add_title(s, "Riesgos, próximos pasos y decisiones",
          eyebrow="Datos provenientes de Jira")

# Riesgos y bloqueos
add_text(s, Inches(0.5), Inches(2.3), Inches(8), Inches(0.35),
         "Riesgos y bloqueos", font=FONT_HEAD, size=14, color=WHITE, bold=True)
add_branded_table(
    s, Inches(0.5), Inches(2.75), Inches(12.3),
    ["Descripción", "Mitigación / acción", "Responsable", "Severidad"], RIESGOS,
    col_fr=[4.0, 4.0, 1.8, 1.5],
)

riesgos_bottom = Inches(2.75) + Inches(0.42) * (len(RIESGOS) + 1) + Inches(0.3)

# Próximos pasos
add_text(s, Inches(0.5), riesgos_bottom, Inches(6), Inches(0.3),
         "PRÓXIMOS PASOS", font=FONT_BODY, size=12, color=LIME, bold=True)
py = riesgos_bottom + Inches(0.4)
for paso in PROXIMOS:
    add_text(s, Inches(0.6), py, Inches(0.3), Inches(0.3), "▸",
             font=FONT_BODY, size=11, color=LIME, bold=True)
    add_text(s, Inches(0.95), py, Inches(6.0), Inches(0.3), paso,
             font=FONT_BODY, size=11, color=WHITE)
    py += Inches(0.32)

# Decisiones / apoyos del cliente
add_text(s, Inches(7.0), riesgos_bottom, Inches(6), Inches(0.3),
         "DECISIONES / APOYOS DEL CLIENTE", font=FONT_BODY, size=12, color=LIME, bold=True)
add_accent_card(s, Inches(7.0), riesgos_bottom + Inches(0.4), Inches(5.8), Inches(1.1),
                fill_color=BG_PANEL)
add_text(s, Inches(7.25), riesgos_bottom + Inches(0.5), Inches(5.4), Inches(0.9),
         DECISIONES, font=FONT_BODY, size=11, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s, 4)


# ═══════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════
import os
os.makedirs(os.path.dirname(OUT) or ".", exist_ok=True)
prs.save(OUT)
print(f"✓ Saved: {OUT}")
print(f"  Slides: {len(prs.slides)}")
